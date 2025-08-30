# app.py
import os
import re
import json
import requests
from io import BytesIO
from flask import Flask, render_template, request, send_file, jsonify
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# ---------------------------
# Config
# ---------------------------
load_dotenv()
app = Flask(__name__, static_folder="static", template_folder="templates")
app.secret_key = os.getenv("FLASK_SECRET", "supersecretkey")

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

EURI_API_KEY = os.getenv("EURI_API_KEY")
if not EURI_API_KEY:
    raise ValueError("⚠️ Missing EURI_API_KEY in .env")

CHAT_URL = "https://api.euron.one/api/v1/euri/chat/completions"
HEADERS = {"Content-Type": "application/json", "Authorization": f"Bearer {EURI_API_KEY}"}

# ---------------------------
# Helpers
# ---------------------------
def extract_json_from_text(text: str):
    """Extract JSON array/object from model text output."""
    if not text:
        raise ValueError("Empty text from model")
    try:
        return json.loads(text)
    except Exception:
        pass
    m = re.search(r'(\[.*\])', text, re.DOTALL)
    if m:
        try:
            return json.loads(m.group(1))
        except Exception:
            pass
    raise ValueError("Could not parse JSON from model output")

def get_slides_from_euri(input_text: str, guidance: str = ""):
    """Call Euron API to generate structured slide data (max 5)."""
    prompt = f"""
Turn the following into a 5-slide PowerPoint presentation.
Return JSON array only, no extra text.

Each slide must have:
- "title": string
- "subtitle": optional string
- "highlights": optional list of 1–3 short phrases
- "bullets": list of concise bullets
- "details": optional speaker notes

Input:
{input_text}

Guidance:
{guidance if guidance else "No special guidance"}
"""
    payload = {"model": "gpt-4.1-nano", "messages": [{"role": "user", "content": prompt}]}
    r = requests.post(CHAT_URL, headers=HEADERS, json=payload, timeout=60)
    r.raise_for_status()
    data = r.json()
    text = data["choices"][0]["message"]["content"]
    return extract_json_from_text(text)[:5]

# ---------------------------
# PPT Utils
# ---------------------------
def delete_slide(prs, slide):
    """Delete a slide safely from a Presentation (fixed version)."""
    slide_id = slide.slide_id
    slides = prs.slides._sldIdLst
    for sldId in list(slides):
        if sldId.attrib["id"] == str(slide_id):
            slides.remove(sldId)
            break
    try:
        slide_part = slide.part
        rId = slide_part._element.get("r:id")
        if rId:
            prs.part.drop_rel(rId)
    except Exception:
        pass

def clear_only_text(prs: Presentation):
    """Clear only text frames, leave design untouched."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                try:
                    shape.text_frame.clear()
                except Exception:
                    pass

def find_best_layout(prs: Presentation):
    """Pick a layout with title+body placeholders if possible."""
    for layout in prs.slide_layouts:
        types = [ph.placeholder_format.type for ph in layout.placeholders if ph.is_placeholder]
        if 0 in types and 1 in types:
            return layout
    return prs.slide_layouts[0]

def safe_set_paragraph_text(paragraph, text, font_size=Pt(18), bold=False, color=RGBColor(0,0,0), level=None):
    paragraph.text = text
    try:
        paragraph.font.size = font_size
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        if level is not None:
            paragraph.level = level
    except Exception:
        pass

def replace_text_in_slide(slide, slide_data: dict):
    """Insert AI text into placeholders, preserving design."""
    title = slide_data.get("title", "")
    highlights = slide_data.get("highlights", []) or []
    bullets = slide_data.get("bullets", []) or []
    details = slide_data.get("details", "")

    # Title
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type == 0 and shape.has_text_frame:
                    shape.text_frame.clear()
                    p = shape.text_frame.add_paragraph()
                    safe_set_paragraph_text(p, title, font_size=Pt(32), bold=True, color=RGBColor(0,51,102))
                    break
            except Exception:
                continue

    # Body
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type == 1 and shape.has_text_frame:
                    shape.text_frame.clear()
                    for hl in highlights:
                        p = shape.text_frame.add_paragraph()
                        safe_set_paragraph_text(p, hl, font_size=Pt(20), bold=True, color=RGBColor(0,102,204))
                    for b in bullets:
                        p = shape.text_frame.add_paragraph()
                        safe_set_paragraph_text(p, b, font_size=Pt(18), color=RGBColor(60,60,60), level=1)
                    break
            except Exception:
                continue

    # Speaker notes
    if details:
        try:
            slide.notes_slide.notes_text_frame.clear()
            slide.notes_slide.notes_text_frame.text = details
        except Exception:
            pass

def create_default_presentation(slides_data: list):
    """Fallback: make a professional 5-slide deck without template."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    bg_colors = [
        RGBColor(10, 25, 64), RGBColor(6, 57, 112),
        RGBColor(58, 2, 86), RGBColor(125, 31, 58), RGBColor(6, 77, 83)
    ]

    for idx, sd in enumerate(slides_data[:5]):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_colors[idx % len(bg_colors)]

        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1.4))
        tf = tb.text_frame
        p = tf.add_paragraph()
        safe_set_paragraph_text(p, sd.get("title", ""), font_size=Pt(38), bold=True, color=RGBColor(255,255,255))

        # Highlights
        y = 2.0
        for hl in sd.get("highlights", []):
            box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.6))
            p = box.text_frame.add_paragraph()
            safe_set_paragraph_text(p, hl, font_size=Pt(22), bold=True, color=RGBColor(255,215,0))
            y += 0.7

        # Bullets
        box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(3))
        for b in sd.get("bullets", []):
            p = box.text_frame.add_paragraph()
            safe_set_paragraph_text(p, b, font_size=Pt(20), color=RGBColor(245,245,245), level=1)

        # Speaker notes
        if sd.get("details"):
            slide.notes_slide.notes_text_frame.text = sd["details"]

    return prs

# ---------------------------
# Flask Routes
# ---------------------------
@app.route("/", methods=["GET"])
def home():
    return render_template("index.html")

@app.route("/", methods=["POST"])
def generate_ppt():
    input_text = request.form.get("input_text", "").strip()
    guidance = request.form.get("guidance", "").strip()
    ppt_file = request.files.get("ppt_template")

    if not input_text:
        return jsonify({"error": "input_text required"}), 400

    try:
        slides_data = get_slides_from_euri(input_text, guidance)
    except Exception as e:
        print("⚠️ AI failed, using fallback:", e)
        slides_data = [
            {"title": "Demo Slide 1", "bullets": ["Fallback bullet 1"]},
            {"title": "Demo Slide 2", "bullets": ["Fallback bullet 2"]},
            {"title": "Demo Slide 3", "bullets": ["Fallback bullet 3"]},
            {"title": "Demo Slide 4", "bullets": ["Fallback bullet 4"]},
            {"title": "Demo Slide 5", "bullets": ["Fallback bullet 5"]},
        ]

    if ppt_file and ppt_file.filename:
        filename = secure_filename(ppt_file.filename)
        path = os.path.join(UPLOAD_FOLDER, filename)
        ppt_file.save(path)
        prs = Presentation(path)

        # Force exactly 5 slides
        while len(prs.slides) > 5:
            delete_slide(prs, prs.slides[-1])
        while len(prs.slides) < 5:
            layout = find_best_layout(prs)
            prs.slides.add_slide(layout)

        # Erase old text
        clear_only_text(prs)

        # Fill with AI data
        for i, sd in enumerate(slides_data):
            replace_text_in_slide(prs.slides[i], sd)
    else:
        prs = create_default_presentation(slides_data)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    print(f"✅ Generated {len(prs.slides)} slides successfully")

    return send_file(
        buf,
        as_attachment=True,
        download_name="generated_presentation.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == "__main__":
     port = int(os.environ.get("PORT", 5000))
     app.run(host="0.0.0.0", port=port)

